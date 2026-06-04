# -*- coding: utf-8 -*-
"""
Générateur de présentation PowerPoint enrichie pour l'oral de clustering.
Génère un diaporama de 9 diapositives au format 16:9 avec une charte Dark Cyber.
Inclut des sections dédiées au clustering physique, au MMSB social et au LDA textuel.
"""

import os
import sys
import subprocess

# Installation automatique de python-pptx si absent
try:
    import pptx
except ImportError:
    print("La bibliothèque 'python-pptx' est absente. Tentative d'installation...")
    try:
        subprocess.check_call([sys.executable, "-m", "pip", "install", "python-pptx"])
        import pptx
        print("Installation réussie de 'python-pptx'.")
    except Exception as e:
        print(f"Erreur d'installation : {e}. Veuillez exécuter 'pip install python-pptx' manuellement.")
        sys.exit(1)

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

# Définition de la palette de couleurs (Thème Dark Cyber)
COLOR_BG = RGBColor(18, 18, 22)        # Fond très sombre (#121216)
COLOR_CARD = RGBColor(30, 30, 38)      # Fond des cartes de contenu (#1E1E26)
COLOR_TEXT_LIGHT = RGBColor(240, 240, 245) # Texte principal (#F0F0F5)
COLOR_TEXT_MUTED = RGBColor(170, 170, 185) # Texte secondaire (#AAAAAB9)

# Accents
COLOR_CYAN = RGBColor(0, 240, 255)     # Accent Cyan Néon
COLOR_PINK = RGBColor(255, 0, 127)     # Accent Rose/Magenta
COLOR_PURPLE = RGBColor(176, 80, 255)   # Accent Violet
COLOR_WHITE = RGBColor(255, 255, 255)

def appliquer_fond_sombre(slide):
    """Ajoute un rectangle couvrant toute la diapositive pour faire office de fond sombre."""
    background = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(7.5)
    )
    background.fill.solid()
    background.fill.fore_color.rgb = COLOR_BG
    background.line.color.rgb = COLOR_BG
    # Envoyer le fond en arrière-plan
    slide.shapes._spTree.remove(background._element)
    slide.shapes._spTree.insert(2, background._element)

def ajouter_titre(slide, texte, accent_color=COLOR_CYAN):
    """Ajoute un titre stylisé avec un liseré coloré vertical sur la gauche (effet premium)."""
    # Liseré vertical
    liseret = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0.6), Inches(0.5), Inches(0.12), Inches(0.8)
    )
    liseret.fill.solid()
    liseret.fill.fore_color.rgb = accent_color
    liseret.line.color.rgb = accent_color
    
    # Texte du titre
    tx_box = slide.shapes.add_textbox(Inches(0.9), Inches(0.4), Inches(11.5), Inches(0.9))
    tf = tx_box.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_top = tf.margin_right = tf.margin_bottom = 0
    
    p = tf.paragraphs[0]
    p.text = texte
    p.font.name = "Trebuchet MS"
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = COLOR_TEXT_LIGHT
    p.alignment = PP_ALIGN.LEFT

def ajouter_carte(slide, x, y, w, h, titre_carte, contenu_bullets, accent_color=COLOR_CYAN):
    """Génère une carte de contenu stylisée (rectangle arrondi sombre avec liseré supérieur)."""
    # Corps de la carte
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    card.fill.solid()
    card.fill.fore_color.rgb = COLOR_CARD
    card.line.color.rgb = COLOR_BG
    
    # Liseré supérieur de la carte
    top_line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(0.08))
    top_line.fill.solid()
    top_line.fill.fore_color.rgb = accent_color
    top_line.line.color.rgb = accent_color
    
    # Zone de texte
    tx_box = slide.shapes.add_textbox(Inches(x + 0.25), Inches(y + 0.2), Inches(w - 0.5), Inches(h - 0.4))
    tf = tx_box.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_top = tf.margin_right = tf.margin_bottom = 0
    
    # Titre de la carte
    p_title = tf.paragraphs[0]
    p_title.text = titre_carte
    p_title.font.name = "Trebuchet MS"
    p_title.font.size = Pt(20)
    p_title.font.bold = True
    p_title.font.color.rgb = accent_color
    p_title.space_after = Pt(12)
    
    # Bullet points
    for bullet in contenu_bullets:
        p = tf.add_paragraph()
        p.text = "•  " + bullet
        p.font.name = "Calibri"
        p.font.size = Pt(14)
        p.font.color.rgb = COLOR_TEXT_LIGHT
        p.space_after = Pt(6)
        p.level = 0

def main():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    blank_layout = prs.slide_layouts[6]
    
    # ==========================================
    # DIAPO 1 : ACCUEIL
    # ==========================================
    slide1 = prs.slides.add_slide(blank_layout)
    appliquer_fond_sombre(slide1)
    
    tx_box = slide1.shapes.add_textbox(Inches(1.0), Inches(2.0), Inches(11.333), Inches(2.5))
    tf = tx_box.text_frame
    tf.word_wrap = True
    
    p = tf.paragraphs[0]
    p.text = "ÉTUDE DE CLUSTERING & GRAPHES"
    p.font.name = "Trebuchet MS"
    p.font.size = Pt(50)
    p.font.bold = True
    p.font.color.rgb = COLOR_CYAN
    
    p2 = tf.add_paragraph()
    p2.text = "Profilage comportemental, réseau social MMSB et analyse textuelle LDA"
    p2.font.name = "Calibri"
    p2.font.size = Pt(22)
    p2.font.color.rgb = COLOR_TEXT_LIGHT
    p2.space_before = Pt(8)
    
    tx_box_author = slide1.shapes.add_textbox(Inches(1.0), Inches(5.2), Inches(6.0), Inches(1.5))
    tf_author = tx_box_author.text_frame
    p_auth = tf_author.paragraphs[0]
    p_auth.text = "Présentation d'oral - Bachelor 2 - ESIEE-IT\nModèles de scikit-learn & algorithme probabiliste MMSB maison\nLivrables : Dataset extrait + Notebook Jupyter commenté"
    p_auth.font.name = "Calibri"
    p_auth.font.size = Pt(13)
    p_auth.font.color.rgb = COLOR_TEXT_MUTED
    
    line = slide1.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1.0), Inches(4.7), Inches(5.5), Inches(0.06))
    line.fill.solid()
    line.fill.fore_color.rgb = COLOR_PINK
    line.line.color.rgb = COLOR_PINK
    
    # ==========================================
    # DIAPO 2 : PROBLEMATIQUE
    # ==========================================
    slide2 = prs.slides.add_slide(blank_layout)
    appliquer_fond_sombre(slide2)
    ajouter_titre(slide2, "Problématique & Cadre Métier", COLOR_CYAN)
    
    ajouter_carte(
        slide2, 0.8, 1.8, 5.5, 4.8, 
        "La Question Métier", 
        [
            "Comment segmenter automatiquement et sans a priori la population de notre Jeu de la Vie ?",
            "Découvrir des profils comportementaux homogènes (reproducteurs actifs, marginaux épuisés, agresseurs).",
            "Permettre aux administrateurs de réguler l'équilibre écologique et social de l'île."
        ], 
        COLOR_CYAN
    )
    
    ajouter_carte(
        slide2, 7.0, 1.8, 5.5, 4.8, 
        "Les Outils IA Non Supervisés", 
        [
            "Clustering Physique : Regrouper selon l'état biologique (énergie, libido, colère, isolement).",
            "Clustering Réseau (MMSB) : Découvrir des groupes sociaux aux appartenances fluides.",
            "Topic Modeling (LDA) : Extraire les thèmes de discussion des agents à partir de leurs dialogues LLM."
        ], 
        COLOR_PINK
    )
    
    # ==========================================
    # DIAPO 3 : COLLECTE & DATASET
    # ==========================================
    slide3 = prs.slides.add_slide(blank_layout)
    appliquer_fond_sombre(slide3)
    ajouter_titre(slide3, "Collecte & Traitement des Données", COLOR_PINK)
    
    ajouter_carte(
        slide3, 0.8, 1.8, 5.5, 4.8, 
        "Extraction par API", 
        [
            "Extraction dynamique via l'API REST de la simulation (FastAPI).",
            "Attributs collectés : Énergie, Libido, Colère, Fatigue sociale, et Nage (couleur de l'agent).",
            "Générateur synthétique de repli extrêmement réaliste (150 échantillons) pour garantir la reproductibilité."
        ], 
        COLOR_PINK
    )
    
    ajouter_carte(
        slide3, 7.0, 1.8, 5.5, 4.8, 
        "Prétraitements & EDA", 
        [
            "Analyse descriptive et visualisations des distributions via Seaborn.",
            "Mise à l'échelle (StandardScaler) : Centrage-réduction des données physiques pour éviter que les échelles ne dominent les calculs de distance.",
            "Création de la matrice d'adjacence sociale pour l'analyse réseau."
        ], 
        COLOR_PURPLE
    )
    
    # ==========================================
    # DIAPO 4 : BENCHMARK MODELES
    # ==========================================
    slide4 = prs.slides.add_slide(blank_layout)
    appliquer_fond_sombre(slide4)
    ajouter_titre(slide4, "Benchmark des Modèles de Clustering", COLOR_PURPLE)
    
    ajouter_carte(
        slide4, 0.8, 1.8, 3.6, 4.8, 
        "K-Means", 
        [
            "Partitionnement dur (classes exclusives).",
            "Recherche de K optimal par la Méthode du Coude (Inertie).",
            "Validation de cohésion spatiale par le score de Silhouette global et individuel."
        ], 
        COLOR_CYAN
    )
    
    ajouter_carte(
        slide4, 4.8, 1.8, 3.6, 4.8, 
        "GMM (Gaussiennes)", 
        [
            "Soft clustering probabiliste.",
            "Sélection rigoureuse de K via les critères BIC et AIC (creux marqué à K=3).",
            "Permet d'évaluer la certitude de l'affectation comportementale d'un agent."
        ], 
        COLOR_PINK
    )
    
    ajouter_carte(
        slide4, 8.8, 1.8, 3.6, 4.8, 
        "DBSCAN & CAH", 
        [
            "DBSCAN : Clustering par densité. Isole les 'marginaux/exclus' dans la classe bruit (-1) sans fixer K.",
            "CAH : Structuration hiérarchique. Construction de dendrogrammes de relations."
        ], 
        COLOR_PURPLE
    )
    
    # ==========================================
    # DIAPO 5 : MMSB (Réseau Social) - NOUVEAU
    # ==========================================
    slide5 = prs.slides.add_slide(blank_layout)
    appliquer_fond_sombre(slide5)
    ajouter_titre(slide5, "Analyse de Réseau : Le Modèle MMSB", COLOR_CYAN)
    
    ajouter_carte(
        slide5, 0.8, 1.8, 5.5, 4.8, 
        "Le Modèle MMSB", 
        [
            "Mixed Membership Stochastic Blockmodel : modèle probabiliste variationnel de réseau social.",
            "Les relations ne sont pas binaires : les agents ont des appartenances communautaires mixtes.",
            "Ajustement sur la matrice d'adjacence des amitiés de l'île."
        ], 
        COLOR_CYAN
    )
    
    ajouter_carte(
        slide5, 7.0, 1.8, 5.5, 4.8, 
        "Résultats de l'Inférence", 
        [
            "Visualisation du graphe (Spring layout) coloré en dégradé continu de Rouge à Bleu.",
            "Identification des agents mixtes jouant le rôle de 'ponts' sociaux.",
            "Matrice de connectivité B : forte assortativité (les membres du même groupe se lient plus volontiers)."
        ], 
        COLOR_PINK
    )
    
    # ==========================================
    # DIAPO 6 : LDA (Topic Modeling) - NOUVEAU
    # ==========================================
    slide6 = prs.slides.add_slide(blank_layout)
    appliquer_fond_sombre(slide6)
    ajouter_titre(slide6, "Topic Modeling sur Dialogues : La LDA", COLOR_PINK)
    
    ajouter_carte(
        slide6, 0.8, 1.8, 5.5, 4.8, 
        "NLP & Extraction", 
        [
            "Traitement du langage naturel sur les journaux de discussion des agents.",
            "Pipeline : Nettoyage, exclusion des mots vides français (stop words) et vectorisation avec CountVectorizer.",
            "Entraînement de la Latent Dirichlet Allocation (LDA) avec K=3 thèmes."
        ], 
        COLOR_PINK
    )
    
    ajouter_carte(
        slide6, 7.0, 1.8, 5.5, 4.8, 
        "Thèmes & Correspondances", 
        [
            "Thème 0 (Amour/Reproduction) : mots comme 'libido', 'promène', 'couple'.",
            "Thème 1 (Repos/Sommeil/Solitude) : mots comme 'épuisé', 'dormir', 'isoler'.",
            "Thème 2 (Colère/Conflit/Menace) : mots comme 'colère', 'frapper', 'haine'.",
            "L'alignement est parfait avec les clusters physiques."
        ], 
        COLOR_PURPLE
    )
    
    # ==========================================
    # DIAPO 7 : PCA & VISUALISATION 2D
    # ==========================================
    slide7 = prs.slides.add_slide(blank_layout)
    appliquer_fond_sombre(slide7)
    ajouter_titre(slide7, "Visualisation 2D (PCA)", COLOR_PURPLE)
    
    ajouter_carte(
        slide7, 0.8, 1.8, 5.5, 4.8, 
        "Rôle de la PCA", 
        [
            "Projection linéaire orthogonale : réduction de 5 variables physiques vers 2 dimensions.",
            "Variance cumulée : plus de 80% de la variance d'origine est conservée sur les 2 premiers axes.",
            "Validation visuelle : les frontières séparant les clusters du GMM sont claires."
        ], 
        COLOR_PURPLE
    )
    
    ajouter_carte(
        slide7, 7.0, 1.8, 5.5, 4.8, 
        "Analyse des Composantes", 
        [
            "Axe 1 (PC1) : Fortement lié à l'énergie (négativement) et à la fatigue sociale (positivement).",
            "Axe 2 (PC2) : Dominé par la colère et l'état d'isolement.",
            "Les zones de chevauchement traduisent la dynamique transitoire de la vie des agents."
        ], 
        COLOR_CYAN
    )
    
    # ==========================================
    # DIAPO 8 : PROFILAGE DES CITOYENS
    # ==========================================
    slide8 = prs.slides.add_slide(blank_layout)
    appliquer_fond_sombre(slide8)
    ajouter_titre(slide8, "Interprétation & Profils Comportementaux", COLOR_CYAN)
    
    ajouter_carte(
        slide8, 0.8, 1.8, 3.6, 4.8, 
        "1. Citoyens Actifs", 
        [
            "Caractéristiques : Énergie ~75, Libido ~80, colère très basse.",
            "Comportement : Force démographique et sociale de l'île. Assurent la reproduction.",
            "Langage : Thème 0 (Amour)."
        ], 
        COLOR_CYAN
    )
    
    ajouter_carte(
        slide8, 4.8, 1.8, 3.6, 4.8, 
        "2. Ermites Isolés", 
        [
            "Caractéristiques : Énergie ~25, Fatigue sociale ~90, libido basse.",
            "Comportement : Retrait social pour se reposer. Se déplacent vers les zones calmes.",
            "Langage : Thème 1 (Repos)."
        ], 
        COLOR_PINK
    )
    
    ajouter_carte(
        slide8, 8.8, 1.8, 3.6, 4.8, 
        "3. Individus Hostiles", 
        [
            "Caractéristiques : Colère ~85, énergie moyenne, libido basse.",
            "Comportement : Agresseurs ou assassins potentiels. Déclenchent des conflits physiques.",
            "Langage : Thème 2 (Conflit)."
        ], 
        COLOR_PURPLE
    )
    
    # ==========================================
    # DIAPO 9 : CONCLUSION
    # ==========================================
    slide9 = prs.slides.add_slide(blank_layout)
    appliquer_fond_sombre(slide9)
    ajouter_titre(slide9, "Conclusion & Perspectives", COLOR_PINK)
    
    ajouter_carte(
        slide9, 0.8, 1.8, 5.5, 4.8, 
        "Résultats Obtenus", 
        [
            "La simulation constitue un excellent bac à sable sociologique : les clusters collent parfaitement aux règles codées.",
            "Benchmark robuste validant les choix des modèles d'IA non supervisés vus en cours.",
            "Livrables prêts : Notebook commenté de A à Z + Dataset + Présentation PPTX."
        ], 
        COLOR_CYAN
    )
    
    ajouter_carte(
        slide9, 7.0, 1.8, 5.5, 4.8, 
        "Limites & Améliorations", 
        [
            "Limites : Caractère purement statique de l'analyse ($t$ unique).",
            "Améliorations : Clustering de séries temporelles (dynamique des jauges dans le temps).",
            "MMSB Temporel : Suivre la dérive des amitiés de l'île à chaque tick."
        ], 
        COLOR_PINK
    )
    
    output_filename = "Jeu_de_la_Vie_Clustering_Presentation.pptx"
    prs.save(output_filename)
    print(f"Présentation PowerPoint mise à jour avec 9 diapositives : {os.path.abspath(output_filename)}")

if __name__ == "__main__":
    main()
